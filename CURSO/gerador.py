from pptx import Presentation

# 1. Cria a apresentação em branco
prs = Presentation()

# Função para facilitar a criação de slides
def add_slide(title_text, body_text):
    slide_layout = prs.slide_layouts[1] # Layout de Título e Conteúdo
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = title_text
    body.text = body_text

# 2. Adicionando os slides com as informações que estudamos
add_slide("Dominando Módulos em Python", 
          "A Cozinha da Programação\n\n- Entendendo Módulos\n- math, random e platform\n- pip, PyPI e .pyd")

add_slide("O que é um Módulo? (A Caixa de Ferramentas)", 
          "- Códigos crescem: dividimos tarefas em módulos.\n"
          "- dir(modulo): O índice da caixa. Mostra tudo que tem dentro.\n"
          "- Formas de importar:\n"
          "  1. import math (Traz a caixa inteira)\n"
          "  2. from math import pi (Traz só uma ferramenta)\n"
          "  3. import math as m (Dá um apelido curto)")

add_slide("Módulo math & Namespaces", 
          "- math: Ferramentas matemáticas de alta precisão.\n"
          "- ceil(x): Arredonda para cima (Teto).\n"
          "- floor(x): Arredonda para baixo (Piso).\n"
          "- Namespaces: Etiquetas de organização. O 'math.pi' não entra em conflito com uma variável 'pi' que você criar.")

add_slide("Módulo random (O Sorteio)", 
          "- Gera números pseudoaleatórios (parecem aleatórios, mas seguem uma matemática).\n"
          "- seed(): A semente. Define de onde o sorteio começa.\n"
          "- choice(): Escolhe um item aleatório de uma lista.\n"
          "- AVISO: Nunca use para senhas! Para isso, use o módulo 'secrets'.")

add_slide("Módulo platform (O Inspetor da Cozinha)", 
          "- Permite que o código saiba onde está rodando.\n"
          "- platform(): Mostra o Sistema Operacional e versão.\n"
          "- machine(): Mostra a arquitetura (ex: x86_64).\n"
          "- processor(): Mostra qual é o cérebro da máquina (ex: Intel Core i5).")

add_slide("pip, PyPI e arquivos .pyd", 
          "- PyPI: O 'Mercado Livre' oficial do Python na internet.\n"
          "- pip: O entregador que baixa e instala pacotes do PyPI.\n"
          "- .pyd: Módulos super-rápidos escritos em C/C++ para Windows, que o Python consegue importar e usar como se fossem nativos.")

# 3. Salva o arquivo final
prs.save('Meus_Estudos_Python.pptx')
print("Sucesso! O arquivo 'Meus_Estudos_Python.pptx' foi gerado na mesma pasta deste script.")